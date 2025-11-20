#!/usr/bin/env python3
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import os

def populate_template(template_path, slides_data, output_path):
    """
    Loads a template PPTX and populates it with AI-generated content
    while preserving ALL design elements
    """
    try:
        # Load the template presentation
        prs = Presentation(template_path)
        
        # Get slide layouts
        slide_layouts = prs.slide_layouts
        
        print(f"Template loaded. Has {len(slide_layouts)} layouts", file=sys.stderr)
        
        # Get number of existing slides
        existing_slide_count = len(prs.slides)
        print(f"Template has {existing_slide_count} existing slides", file=sys.stderr)
        
        # Delete all existing content slides (keep masters/layouts)
        # Work backwards to avoid index issues
        for i in range(existing_slide_count - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
        
        print(f"Cleared existing slides. Creating {len(slides_data)} new slides", file=sys.stderr)
        
        # Create new slides based on AI content
        for idx, slide_content in enumerate(slides_data):
            print(f"Creating slide {idx + 1}: {slide_content.get('title', 'Untitled')}", file=sys.stderr)
            
            # Determine layout based on slide type
            layout = get_best_layout(slide_layouts, slide_content.get('type', 'content'))
            
            # Add slide with the layout
            slide = prs.slides.add_slide(layout)
            
            # Populate slide with content
            populate_slide_content(slide, slide_content)
        
        # Save the presentation
        prs.save(output_path)
        
        print(f"Presentation saved to {output_path}", file=sys.stderr)
        
        return {
            'success': True,
            'output_path': output_path,
            'slide_count': len(slides_data)
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc(file=sys.stderr)
        return {
            'success': False,
            'error': str(e)
        }

def get_best_layout(slide_layouts, slide_type):
    """
    Selects the best slide layout based on slide type
    """
    # Try to find layout by name or index
    if slide_type == 'title':
        # Look for title slide layout (usually index 0)
        for layout in slide_layouts:
            if 'title' in layout.name.lower() and 'only' not in layout.name.lower():
                return layout
        return slide_layouts[0]  # Default to first layout
    else:
        # Look for content layout (usually index 1 or has "content" in name)
        for layout in slide_layouts:
            if any(word in layout.name.lower() for word in ['content', 'bullet', 'text']):
                return layout
        # Default to second layout if available, otherwise first
        return slide_layouts[1] if len(slide_layouts) > 1 else slide_layouts[0]

def populate_slide_content(slide, content):
    """
    Populates a slide with content while preserving template formatting
    """
    title_text = content.get('title', '')
    main_content = content.get('content', '')
    bullet_points = content.get('bulletPoints', [])
    
    print(f"  Populating: Title='{title_text}', Bullets={len(bullet_points)}", file=sys.stderr)
    
    # Strategy: Find and fill placeholders
    title_filled = False
    content_filled = False
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        # Check if it's a placeholder
        if shape.is_placeholder:
            phf = shape.placeholder_format
            
            # Title placeholder (type 1)
            if phf.type == 1 and not title_filled and title_text:
                shape.text = title_text
                title_filled = True
                print(f"    Set title placeholder", file=sys.stderr)
            
            # Content/Body placeholder (type 2)
            elif phf.type == 2 and not content_filled:
                fill_content_placeholder(shape, main_content, bullet_points)
                content_filled = True
                print(f"    Set content placeholder", file=sys.stderr)
            
            # Subtitle placeholder (type 3) - use for content on title slides
            elif phf.type == 3 and not content_filled and main_content:
                shape.text = main_content
                content_filled = True
                print(f"    Set subtitle placeholder", file=sys.stderr)
    
    # Fallback: If no placeholders found, try to find text boxes
    if not title_filled and title_text:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == '':
                shape.text = title_text
                title_filled = True
                break
    
    if not content_filled and (main_content or bullet_points):
        for shape in slide.shapes:
            if shape.has_text_frame and not title_filled:
                fill_content_placeholder(shape, main_content, bullet_points)
                content_filled = True
                break

def fill_content_placeholder(shape, main_content, bullet_points):
    """
    Fills a content placeholder with text and/or bullet points
    """
    text_frame = shape.text_frame
    text_frame.clear()
    
    # If there are bullet points, add them
    if bullet_points and len(bullet_points) > 0:
        for i, bullet in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            # Preserve existing formatting if possible
    
    # If there's main content but no bullets, add it
    elif main_content:
        p = text_frame.paragraphs[0]
        p.text = main_content

def main():
    if len(sys.argv) != 4:
        result = {
            'success': False,
            'error': 'Usage: template_populator.py <template_path> <slides_json> <output_path>'
        }
        print(json.dumps(result))
        sys.exit(1)
    
    template_path = sys.argv[1]
    slides_json = sys.argv[2]
    output_path = sys.argv[3]
    
    # Verify template exists
    if not os.path.exists(template_path):
        result = {
            'success': False,
            'error': f'Template file not found: {template_path}'
        }
        print(json.dumps(result))
        sys.exit(1)
    
    try:
        slides_data = json.loads(slides_json)
        result = populate_template(template_path, slides_data, output_path)
        print(json.dumps(result))
        
        if not result['success']:
            sys.exit(1)
            
    except Exception as e:
        import traceback
        traceback.print_exc(file=sys.stderr)
        result = {
            'success': False,
            'error': str(e)
        }
        print(json.dumps(result))
        sys.exit(1)

if __name__ == '__main__':
    main()

